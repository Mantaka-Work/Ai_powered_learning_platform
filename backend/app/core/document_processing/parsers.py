"""Document parsers for various file types."""
from typing import List, Dict, Any, Optional
from pathlib import Path
import io


class DocumentParser:
    """Parse various document types to extract text content."""
    
    async def parse(self, file_content: bytes, filename: str) -> str:
        """
        Parse document and extract text content.
        
        Args:
            file_content: Raw file bytes
            filename: Original filename (used to determine type)
        
        Returns:
            Extracted text content
        """
        ext = Path(filename).suffix.lower()
        
        if ext == ".pdf":
            return await self._parse_pdf(file_content)
        elif ext in [".pptx", ".ppt"]:
            return await self._parse_pptx(file_content)
        elif ext in [".docx", ".doc"]:
            return await self._parse_docx(file_content)
        elif ext in [".md", ".txt", ".json", ".yaml", ".yml"]:
            return await self._parse_text(file_content)
        elif ext in [".py", ".js", ".ts", ".java", ".cpp", ".c", ".cs", ".go", ".rs"]:
            return await self._parse_code(file_content, ext)
        else:
            # Try as plain text
            return await self._parse_text(file_content)
    
    async def _parse_pdf(self, content: bytes) -> str:
        """Parse PDF file."""
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(io.BytesIO(content))
            text_parts = []
            
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to parse PDF: {str(e)}")
    
    async def _parse_pptx(self, content: bytes) -> str:
        """Parse PowerPoint file."""
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to parse PPTX: {str(e)}")
    
    async def _parse_docx(self, content: bytes) -> str:
        """Parse Word document."""
        try:
            from docx import Document
            
            doc = Document(io.BytesIO(content))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells if cell.text.strip()]
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to parse DOCX: {str(e)}")
    
    async def _parse_text(self, content: bytes) -> str:
        """Parse plain text file."""
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return content.decode("latin-1")
            except Exception as e:
                raise ValueError(f"Failed to decode text: {str(e)}")
    
    async def _parse_code(self, content: bytes, ext: str) -> str:
        """Parse code file with language annotation."""
        text = await self._parse_text(content)
        
        # Map extensions to language names
        lang_map = {
            ".py": "python",
            ".js": "javascript",
            ".ts": "typescript",
            ".java": "java",
            ".cpp": "cpp",
            ".c": "c",
            ".cs": "csharp",
            ".go": "go",
            ".rs": "rust"
        }
        
        lang = lang_map.get(ext, "code")
        
        # Wrap in markdown code block for better processing
        return f"```{lang}\n{text}\n```"


# Singleton instance
_parser = None


def get_parser() -> DocumentParser:
    """Get or create parser singleton."""
    global _parser
    if _parser is None:
        _parser = DocumentParser()
    return _parser
